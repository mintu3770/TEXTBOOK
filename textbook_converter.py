# textbook_converter.py
import os
import re
import tempfile
import base64
import sys
import subprocess
from pathlib import Path
import streamlit as st

# Check and install required packages
required_packages = [
    "python-pptx", 
    "pymupdf", 
    "reportlab", 
    "streamlit",
    "Pillow"  # For image processing
]

def install_packages():
    for package in required_packages:
        try:
            __import__(package.split('-')[0])  # Handle hyphenated names
        except ImportError:
            st.warning(f"Installing missing package: {package}")
            subprocess.check_call([sys.executable, "-m", "pip", "install", package])
            
install_packages()

# Now import the main libraries
try:
    from pptx import Presentation
    import fitz  # PyMuPDF
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import (
        SimpleDocTemplate, Paragraph, Spacer, PageBreak, 
        Image, Frame, PageTemplate
    )
    from reportlab.lib.units import inch
    from PIL import Image as PILImage  # For image validation
except ImportError as e:
    st.error(f"Critical import failed: {str(e)}")
    st.stop()

# Set page config
st.set_page_config(
    page_title="Textbook Converter",
    page_icon="📚",
    layout="centered",
    initial_sidebar_state="expanded"
)

# =====================
# TEXT PROCESSING UTILS
# =====================
def clean_text(text):
    """Normalize and clean extracted text"""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text)  # Remove extra whitespace
    text = re.sub(r'•\s*', '\n• ', text)  # Format bullet points
    text = re.sub(r'\b(\d+\.)\s+', r'\n\1 ', text)  # Numbered lists
    text = re.sub(r'([.!?])([A-Z])', r'\1\n\2', text)  # Sentence breaks
    return text.strip()

def create_styles():
    """Create textbook-like paragraph styles"""
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        name='Textbook',
        fontSize=10,
        leading=14,
        spaceAfter=6,
        fontName='Times-Roman',
        alignment=4  # Justified
    ))
    styles.add(ParagraphStyle(
        name='Heading1',
        fontSize=16,
        leading=18,
        spaceAfter=12,
        fontName='Times-Bold'
    ))
    styles.add(ParagraphStyle(
        name='Heading2',
        fontSize=14,
        leading=16,
        spaceAfter=8,
        fontName='Times-BoldItalic'
    ))
    styles.add(ParagraphStyle(
        name='Caption',
        fontSize=9,
        leading=11,
        spaceAfter=12,
        fontName='Times-Italic',
        alignment=1  # Centered
    ))
    return styles

# =====================
# CONTENT EXTRACTION
# =====================
def extract_pptx(pptx_path, img_dir):
    """Extract text and images from PowerPoint files"""
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        st.error(f"Failed to open PPTX file: {str(e)}")
        return []
    
    content = []
    
    for slide_number, slide in enumerate(prs.slides):
        slide_content = {"text": "", "images": []}
        
        # Extract text
        text_elements = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = ' '.join(run.text for run in paragraph.runs if run.text.strip())
                    if text:
                        text_elements.append(text)
        
        slide_content["text"] = clean_text("\n".join(text_elements))
        
        # Extract images
        img_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "image") and hasattr(shape.image, "blob"):
                try:
                    img = shape.image
                    img_bytes = img.blob
                    img_ext = img.ext
                    img_path = img_dir / f"slide_{slide_number}_img_{img_count}.{img_ext}"
                    
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    
                    # Validate image
                    try:
                        PILImage.open(img_path)  # Verify it's a valid image
                        slide_content["images"].append(str(img_path))
                        img_count += 1
                    except:
                        st.warning(f"Skipped invalid image in slide {slide_number+1}")
                        os.remove(img_path)
                except Exception as e:
                    st.warning(f"Failed to extract image: {str(e)}")
        
        content.append(slide_content)
    
    return content

def extract_pdf(pdf_path, img_dir):
    """Extract text and images from PDF files"""
    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        st.error(f"Failed to open PDF file: {str(e)}")
        return []
    
    content = []
    
    for page_number in range(len(doc)):
        page = doc.load_page(page_number)
        page_content = {"text": "", "images": []}
        
        # Extract text
        try:
            text = page.get_text()
            page_content["text"] = clean_text(text)
        except:
            st.warning(f"Text extraction failed on page {page_number+1}")
        
        # Extract images
        try:
            img_list = page.get_images(full=True)
            for img_index, img in enumerate(img_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    img_bytes = base_image["image"]
                    img_ext = base_image["ext"]
                    
                    # Handle JPEG2000 format
                    if img_ext == "jp2":
                        img_ext = "jpeg"
                    
                    img_path = img_dir / f"page_{page_number}_img_{img_index}.{img_ext}"
                    
                    with open(img_path, "wb") as f:
                        f.write(img_bytes)
                    
                    # Validate image
                    try:
                        PILImage.open(img_path)
                        page_content["images"].append(str(img_path))
                    except:
                        st.warning(f"Skipped invalid image on page {page_number+1}")
                        os.remove(img_path)
                except Exception as e:
                    st.warning(f"Image extraction failed: {str(e)}")
        except:
            st.warning(f"Image extraction failed on page {page_number+1}")
        
        content.append(page_content)
    
    return content

# =====================
# PDF GENERATION (WITH ERROR HANDLING)
# =====================
def create_textbook_pdf(content, output_path):
    """Generate textbook-style PDF with two-column layout"""
    if not content:
        st.error("No content to generate PDF")
        return None
    
    try:
        styles = create_styles()
    except:
        st.error("Failed to create PDF styles")
        return None
    
    try:
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=36,
            leftMargin=36,
            topMargin=36,
            bottomMargin=36
        )
    except Exception as e:
        st.error(f"PDF setup failed: {str(e)}")
        return None
    
    # Two-column layout
    try:
        frame1 = Frame(
            doc.leftMargin, 
            doc.bottomMargin, 
            doc.width/2 - 6, 
            doc.height,
            leftPadding=0,
            bottomPadding=0,
            rightPadding=12,
            topPadding=0
        )
        frame2 = Frame(
            doc.leftMargin + doc.width/2 + 6, 
            doc.bottomMargin, 
            doc.width/2 - 6, 
            doc.height,
            leftPadding=12,
            bottomPadding=0,
            rightPadding=0,
            topPadding=0
        )
        
        doc.addPageTemplates([PageTemplate(id='TwoCol', frames=[frame1, frame2])])
    except:
        st.error("Failed to create page layout")
        return None
    
    elements = []
    
    try:
        # Add title page
        elements.append(Paragraph("Course Materials Textbook", styles['Heading1']))
        elements.append(Spacer(1, 0.5*inch))
        elements.append(Paragraph("Generated from Lecture Materials", styles['Textbook']))
        elements.append(PageBreak())
        
        # Add content
        for i, item in enumerate(content):
            if i > 0:
                elements.append(PageBreak())
            
            # Add text content if exists
            if item.get('text'):
                # Add headings
                if "heading" in item['text'].lower()[:20]:
                    elements.append(Paragraph(item['text'], styles['Heading2']))
                    elements.append(Spacer(1, 0.1*inch))
                else:
                    text_paragraphs = item['text'].split('\n')
                    for p in text_paragraphs:
                        if p.strip():
                            elements.append(Paragraph(p, styles['Textbook']))
            
            # Add images if exists
            if item.get('images'):
                for img_path in item['images']:
                    elements.append(Spacer(1, 0.1*inch))
                    try:
                        # Auto-resize while maintaining aspect ratio
                        img = PILImage.open(img_path)
                        width, height = img.size
                        aspect = width / height
                        
                        # Determine optimal size for textbook
                        max_width = 3 * inch
                        max_height = 2.5 * inch
                        
                        if width > max_width:
                            height = max_width / aspect
                            width = max_width
                        if height > max_height:
                            width = max_height * aspect
                            height = max_height
                            
                        img.close()
                        
                        # Add to PDF
                        pdf_img = Image(img_path, width=width, height=height)
                        pdf_img.hAlign = 'CENTER'
                        elements.append(pdf_img)
                        elements.append(Paragraph(
                            f"Figure {i+1}: Relevant diagram", 
                            styles['Caption']
                        ))
                        elements.append(Spacer(1, 0.2*inch))
                    except Exception as e:
                        st.warning(f"Couldn't process image: {img_path} - {str(e)}")
        
        doc.build(elements)
        return output_path
    except Exception as e:
        st.error(f"PDF generation failed: {str(e)}")
        return None

# =====================
# STREAMLIT UI WITH ENHANCED FEATURES
# =====================
def main():
    st.title("📚 Lecture to Textbook Converter")
    st.markdown("""
    Convert your lecture slides (PPTX/PDF) into textbook-style PDFs for open-book exams.
    """)
    
    with st.sidebar:
        st.header("How to Use")
        st.markdown("""
        1. Upload PPTX or PDF lecture file
        2. Click 'Convert to Textbook'
        3. Download your formatted PDF
        """)
        st.markdown("---")
        st.info("""
        **Features:**
        - Preserves all text content
        - Extracts and resizes images
        - Creates professional two-column layout
        - Automatic heading detection
        """)
        st.warning("""
        **Limitations:**
        - Complex layouts may not convert perfectly
        - Files >50MB may take longer to process
        - Equations remain as images
        """)
    
    uploaded_file = st.file_uploader(
        "Upload lecture file (PPTX or PDF)", 
        type=["pptx", "pdf"],
        accept_multiple_files=False
    )
    
    if uploaded_file is not None:
        # Check file size
        if uploaded_file.size > 50 * 1024 * 1024:  # 50MB limit
            st.error("File size exceeds 50MB limit. Please upload a smaller file.")
            return
            
        with st.spinner("Processing your file..."):
            # Create temp directory
            with tempfile.TemporaryDirectory() as tmp_dir:
                tmp_path = Path(tmp_dir)
                input_path = tmp_path / uploaded_file.name
                
                # Save uploaded file
                try:
                    with open(input_path, "wb") as f:
                        f.write(uploaded_file.getbuffer())
                except Exception as e:
                    st.error(f"Failed to save file: {str(e)}")
                    return
                
                # Extract content
                img_dir = tmp_path / "images"
                img_dir.mkdir(exist_ok=True)
                
                try:
                    if uploaded_file.name.lower().endswith('.pptx'):
                        content = extract_pptx(input_path, img_dir)
                    elif uploaded_file.name.lower().endswith('.pdf'):
                        content = extract_pdf(input_path, img_dir)
                    else:
                        st.error("Unsupported file format")
                        return
                    
                    # Check if we got any content
                    if not content or all(len(item.get('text', '')) == 0 and len(item.get('images', [])) == 0 for item in content):
                        st.error("No extractable content found in the file")
                        return
                    
                    # Generate textbook
                    output_path = tmp_path / "textbook_output.pdf"
                    if create_textbook_pdf(content, output_path):
                        # Show success message
                        st.success("✅ Textbook generated successfully!")
                        
                        # Show preview
                        st.subheader("Preview")
                        try:
                            with open(output_path, "rb") as f:
                                base64_pdf = base64.b64encode(f.read()).decode('utf-8')
                                pdf_display = f'<iframe src="data:application/pdf;base64,{base64_pdf}" width="100%" height="800" type="application/pdf"></iframe>'
                                st.markdown(pdf_display, unsafe_allow_html=True)
                        except:
                            st.warning("Preview unavailable. Please download the file.")
                        
                        # Download button
                        st.download_button(
                            label="Download Textbook PDF",
                            data=open(output_path, "rb").read(),
                            file_name="textbook_output.pdf",
                            mime="application/pdf"
                        )
                    
                except Exception as e:
                    st.error(f"Error processing file: {str(e)}")
                    st.exception(e)

if __name__ == "__main__":
    main()
